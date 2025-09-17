import json
import boto3
import os
import tempfile
from pptx import Presentation

s3 = boto3.client("s3")
BUCKET_NAME = "n8n-ankideckbuilder-pptx-bucket"

def lambda_handler(event, context):
    """
    Expected input (from n8n):
    {
        "s3_key": "myfolder/myslides.pptx"
    }
    """

    try:
        s3_key = event.get("s3_key")
        if not s3_key:
            return {
                "statusCode": 400,
                "body": json.dumps({"error": "Missing 's3_key' in input"})
            }

        # Download file from S3
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            s3.download_file(BUCKET_NAME, s3_key, tmp.name)
            tmp_path = tmp.name

        # Parse PPTX
        prs = Presentation(tmp_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            slides_text.append({
                "slide_number": slide_num,
                "text": "\n".join(text_parts)
            })

        return {
            "statusCode": 200,
            "body": json.dumps({"slides": slides_text})
        }

    except Exception as e:
        return {
            "statusCode": 500,
            "body": json.dumps({"error": str(e)})
        }